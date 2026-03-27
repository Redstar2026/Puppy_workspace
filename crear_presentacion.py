"""Genera presentacion ejecutiva Git & GitHub en formato PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta Walmart ────────────────────────────────────────────────
WALMART_BLUE   = RGBColor(0x00, 0x53, 0xE2)   # #0053e2
WALMART_YELLOW = RGBColor(0xFF, 0xC2, 0x20)   # #ffc220
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BLUE      = RGBColor(0x00, 0x2D, 0x8A)   # azul oscuro
LIGHT_GRAY     = RGBColor(0xF5, 0xF6, 0xFA)
DARK_GRAY      = RGBColor(0x33, 0x33, 0x33)
GREEN          = RGBColor(0x2A, 0x87, 0x03)
RED            = RGBColor(0xEA, 0x11, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if transparency:
        shape.fill.fore_color.theme_color = None
    return shape


def add_text_box(slide, text, x, y, w, h, font_size=18, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_bullet_box(slide, items, x, y, w, h, font_size=16,
                   color=DARK_GRAY, bullet="•"):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def bg_gradient(slide):
    """Fondo degradado azul Walmart."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WALMART_BLUE)
    accent = slide.shapes.add_shape(1, Inches(9), 0, Inches(4.33), SLIDE_H)
    accent.fill.solid()
    accent.fill.fore_color.rgb = DARK_BLUE
    accent.line.fill.background()


def accent_bar(slide, color=WALMART_YELLOW, height=Inches(0.08)):
    """Barra amarilla decorativa en la parte superior."""
    add_rect(slide, 0, 0, SLIDE_W, height, color)


def slide_header(slide, title, subtitle=None):
    """Header estandar para diapositivas internas."""
    accent_bar(slide)
    add_rect(slide, 0, Inches(0.08), SLIDE_W, Inches(1.1), WALMART_BLUE)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.1), Inches(11), Inches(0.9),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.9), Inches(11), Inches(0.4),
                     font_size=14, color=WALMART_YELLOW)


# ── Diapositivas ──────────────────────────────────────────────────
def slide_portada(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg_gradient(slide)

    # Logo text Walmart
    add_rect(slide, Inches(0.4), Inches(0.4), Inches(2.2), Inches(0.55),
             WALMART_YELLOW)
    add_text_box(slide, "  Walmart Tech",
                 Inches(0.4), Inches(0.38), Inches(2.5), Inches(0.6),
                 font_size=16, bold=True, color=DARK_BLUE, align=PP_ALIGN.LEFT)

    # Spark icon simulado
    add_text_box(slide, "✦",
                 Inches(0.4), Inches(1.3), Inches(0.6), Inches(0.6),
                 font_size=32, color=WALMART_YELLOW)

    # Titulo principal
    add_text_box(slide, "Git & GitHub",
                 Inches(0.4), Inches(1.6), Inches(8), Inches(1.4),
                 font_size=54, bold=True, color=WHITE)
    add_text_box(slide, "para el Equipo",
                 Inches(0.4), Inches(2.8), Inches(8), Inches(1.0),
                 font_size=40, bold=False, color=WALMART_YELLOW)
    add_text_box(slide,
                 "Cómo guardar, versionar y compartir\nnuestro código con confianza",
                 Inches(0.4), Inches(3.9), Inches(7.5), Inches(1.0),
                 font_size=18, color=WHITE)

    # Tags
    for i, tag in enumerate(["  Git  ", "  GitHub  ", "  Code Puppy 🐶  "]):
        add_rect(slide,
                 Inches(0.4 + i * 2.1), Inches(5.2),
                 Inches(1.9), Inches(0.45),
                 DARK_BLUE)
        add_text_box(slide, tag,
                     Inches(0.4 + i * 2.1), Inches(5.19),
                     Inches(1.9), Inches(0.45),
                     font_size=14, color=WALMART_YELLOW,
                     align=PP_ALIGN.CENTER)

    # Fecha
    add_text_box(slide, "Walmart · 2026",
                 Inches(0.4), Inches(6.8), Inches(4), Inches(0.4),
                 font_size=12, color=WHITE)


def slide_por_que(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "¿Por qué lo necesitamos?",
                 "El problema que resuelve Git en nuestro equipo")

    # Columna PROBLEMA
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(5.5), WHITE)
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.55), RED)
    add_text_box(slide, "  ❌  Sin Git (hoy)",
                 Inches(0.4), Inches(1.5), Inches(5.8), Inches(0.55),
                 font_size=17, bold=True, color=WHITE)
    add_bullet_box(slide,
                   ["Los scripts se pierden o se sobreescriben",
                    "No sabemos quién cambió qué ni cuándo",
                    "Cada quien tiene su propia versión del archivo",
                    "Imposible volver a una versión anterior",
                    "Compartir archivos por correo o Teams 😅"],
                   Inches(0.6), Inches(2.2), Inches(5.4), Inches(4.0),
                   font_size=15, color=DARK_GRAY)

    # Columna SOLUCIÓN
    add_rect(slide, Inches(6.7), Inches(1.5), Inches(6.2), Inches(5.5), WHITE)
    add_rect(slide, Inches(6.7), Inches(1.5), Inches(6.2), Inches(0.55), GREEN)
    add_text_box(slide, "  ✅  Con Git + GitHub",
                 Inches(6.7), Inches(1.5), Inches(6.2), Inches(0.55),
                 font_size=17, bold=True, color=WHITE)
    add_bullet_box(slide,
                   ["Historial completo de todos los cambios",
                    "Sabemos exactamente qué cambió y quién",
                    "El equipo trabaja sobre la misma versión",
                    "Volvemos a cualquier punto en el tiempo",
                    "Colaboración en tiempo real sin conflictos"],
                   Inches(6.9), Inches(2.2), Inches(5.8), Inches(4.0),
                   font_size=15, color=DARK_GRAY)

    # Flecha central
    add_text_box(slide, "→",
                 Inches(6.1), Inches(3.5), Inches(0.7), Inches(0.7),
                 font_size=36, bold=True, color=WALMART_BLUE,
                 align=PP_ALIGN.CENTER)


def slide_que_es(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "¿Qué es Git y GitHub?",
                 "Las tres herramientas de nuestro flujo de trabajo")

    cards = [
        ("🕰️", "Git", "Máquina del tiempo para tu código",
         ["Guarda cada versión de tus archivos",
          "Vive en tu computadora (local)",
          "Rastrea cada cambio con fecha y autor",
          "Gratis y open source"],
         WALMART_BLUE),
        ("☁️", "GitHub", "La nube donde vive tu código",
         ["Repositorio remoto en internet",
          "Comparte código con tu equipo",
          "Interfaz web para ver cambios",
          "Backup automático en la nube"],
         DARK_BLUE),
        ("🐶", "Code Puppy", "Tu asistente de IA",
         ["Crea el código por ti con IA",
          "Guarda en GitHub automáticamente",
          "Solo dile: 'guarda los cambios'",
          "Disponible en Walmart VPN"],
         RGBColor(0x00, 0x6B, 0xFF)),
    ]

    for i, (icon, title, subtitle, bullets, color) in enumerate(cards):
        x = Inches(0.3 + i * 4.35)
        add_rect(slide, x, Inches(1.6), Inches(4.1), Inches(5.5), WHITE)
        add_rect(slide, x, Inches(1.6), Inches(4.1), Inches(1.2), color)
        add_text_box(slide, icon,
                     x, Inches(1.65), Inches(4.1), Inches(0.6),
                     font_size=28, align=PP_ALIGN.CENTER, color=WHITE)
        add_text_box(slide, title,
                     x, Inches(2.1), Inches(4.1), Inches(0.5),
                     font_size=20, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, subtitle,
                     x + Inches(0.1), Inches(2.85), Inches(3.9), Inches(0.5),
                     font_size=12, color=color, align=PP_ALIGN.CENTER)
        add_bullet_box(slide, bullets,
                       x + Inches(0.15), Inches(3.35), Inches(3.8), Inches(3.5),
                       font_size=13, color=DARK_GRAY)


def slide_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "Cómo empezar: 4 pasos simples",
                 "Una sola vez por persona — luego es automático")

    pasos = [
        ("1", "Crear cuenta", "github.com",
         "Regístrate gratis en github.com con tu correo de Walmart o personal."),
        ("2", "Crear repositorio", "Tu carpeta en la nube",
         "Crea un repo privado en GitHub. Es como una carpeta especial en internet."),
        ("3", "Conectar local", "git init + git remote",
         "Le dices a Git dónde está tu carpeta y la enlazas con GitHub."),
        ("4", "¡Guardar!", "add → commit → push",
         "Cada vez que crees o modifiques algo, lo guardas con 3 comandos."),
    ]

    for i, (num, title, subtitle, desc) in enumerate(pasos):
        x = Inches(0.25 + i * 3.28)
        # Tarjeta
        add_rect(slide, x, Inches(1.6), Inches(3.1), Inches(5.5), WHITE)
        # Número círculo
        add_rect(slide, x + Inches(0.95), Inches(1.75),
                 Inches(1.2), Inches(1.0), WALMART_BLUE)
        add_text_box(slide, num,
                     x + Inches(0.95), Inches(1.8),
                     Inches(1.2), Inches(0.85),
                     font_size=32, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, title,
                     x, Inches(2.85), Inches(3.1), Inches(0.55),
                     font_size=17, bold=True, color=DARK_GRAY,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, subtitle,
                     x, Inches(3.35), Inches(3.1), Inches(0.45),
                     font_size=12, color=WALMART_BLUE,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, desc,
                     x + Inches(0.1), Inches(3.85), Inches(2.9), Inches(2.8),
                     font_size=13, color=DARK_GRAY, align=PP_ALIGN.CENTER)
        # Flecha entre pasos
        if i < 3:
            add_text_box(slide, "→",
                         x + Inches(3.1), Inches(2.9),
                         Inches(0.18), Inches(0.6),
                         font_size=22, bold=True, color=WALMART_BLUE)


def slide_flujo_diario(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "Flujo diario del equipo",
                 "3 comandos que usarás todos los días")

    # Diagrama de flujo
    steps = [
        ("💻", "Tu código\nlocal", WALMART_BLUE),
        ("📋", "git add .", DARK_BLUE),
        ("💾", 'git commit\n-m "mensaje"', DARK_BLUE),
        ("🚀", "git push", DARK_BLUE),
        ("☁️", "GitHub", GREEN),
    ]

    for i, (icon, label, color) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        add_rect(slide, x, Inches(1.8), Inches(2.2), Inches(1.6), color)
        add_text_box(slide, icon,
                     x, Inches(1.85), Inches(2.2), Inches(0.6),
                     font_size=24, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label,
                     x, Inches(2.4), Inches(2.2), Inches(0.85),
                     font_size=13, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER)
        if i < 4:
            add_text_box(slide, "→",
                         x + Inches(2.2), Inches(2.3),
                         Inches(0.35), Inches(0.6),
                         font_size=24, bold=True, color=WALMART_BLUE)

    # Atajo mágico
    add_rect(slide, Inches(0.4), Inches(3.8), Inches(12.5), Inches(1.3),
             DARK_BLUE)
    add_text_box(slide, "⚡  El atajo en una sola línea:",
                 Inches(0.6), Inches(3.85), Inches(12), Inches(0.45),
                 font_size=14, color=WALMART_YELLOW, bold=True)
    add_text_box(slide,
                 "git add .  &&  git commit -m \"descripción\"  &&  git push",
                 Inches(0.6), Inches(4.25), Inches(12), Inches(0.6),
                 font_size=16, bold=True, color=WHITE)

    # Code Puppy tip
    add_rect(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.8), WHITE)
    add_rect(slide, Inches(0.4), Inches(5.3), Inches(0.12), Inches(1.8),
             WALMART_YELLOW)
    add_text_box(slide, "🐶  Tip Code Puppy",
                 Inches(0.65), Inches(5.35), Inches(12), Inches(0.5),
                 font_size=15, bold=True, color=WALMART_BLUE)
    add_text_box(slide,
                 'Solo dile a Code Puppy: "Guarda los cambios en GitHub"'
                 " — y él hace el add, commit y push automáticamente por ti.",
                 Inches(0.65), Inches(5.8), Inches(12), Inches(1.0),
                 font_size=14, color=DARK_GRAY)


def slide_seguridad(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "Reglas de seguridad",
                 "Qué subimos y qué no — política InfoSec Walmart")

    # SI subir
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(5.9), Inches(5.7), WHITE)
    add_rect(slide, Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.6), GREEN)
    add_text_box(slide, "  ✅  SÍ subir a GitHub",
                 Inches(0.4), Inches(1.5), Inches(5.9), Inches(0.6),
                 font_size=17, bold=True, color=WHITE)
    add_bullet_box(slide,
                   ["Scripts Python (.py)",
                    "Consultas SQL (.sql)",
                    "Documentación y guías (.md)",
                    "Presentaciones generadas (.html)",
                    "Código sin datos personales",
                    "Archivos de configuración (.gitignore)"],
                   Inches(0.6), Inches(2.2), Inches(5.5), Inches(4.7),
                   font_size=14, color=DARK_GRAY, bullet="✔")

    # NO subir
    add_rect(slide, Inches(7.0), Inches(1.5), Inches(5.9), Inches(5.7), WHITE)
    add_rect(slide, Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.6), RED)
    add_text_box(slide, "  ❌  NO subir a GitHub",
                 Inches(7.0), Inches(1.5), Inches(5.9), Inches(0.6),
                 font_size=17, bold=True, color=WHITE)
    add_bullet_box(slide,
                   ["Datos con PII (CURP, RFC, etc.)",
                    "Datos de pacientes (HIPAA)",
                    "Archivos CSV / Excel con info sensible",
                    "JSONs con resultados masivos de BQ",
                    "Tokens o contraseñas (.env)",
                    "Bases de datos SQLite"],
                   Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.7),
                   font_size=14, color=DARK_GRAY, bullet="✖")

    # VS
    add_text_box(slide, "VS",
                 Inches(6.2), Inches(3.6), Inches(0.8), Inches(0.7),
                 font_size=22, bold=True, color=WALMART_BLUE,
                 align=PP_ALIGN.CENTER)


def slide_beneficios(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT_GRAY)
    slide_header(slide, "Beneficios para el equipo",
                 "Lo que ganamos al adoptar Git y GitHub")

    benefits = [
        ("🕰️", "Historial Completo",
         "Cada cambio queda registrado con fecha, hora y quién lo hizo."),
        ("🤝", "Colaboración Sin Fricción",
         "Varias personas trabajan en el mismo proyecto sin pisarse."),
        ("🔄", "Recuperación Inmediata",
         "Si algo sale mal, volvemos a cualquier versión anterior en segundos."),
        ("📤", "Compartir Fácilmente",
         "Comparte un link de GitHub en vez de adjuntar archivos por correo."),
        ("🤖", "Automatización con IA",
         "Code Puppy guarda y documenta el código automáticamente."),
        ("🔒", "Código Seguro",
         "Repositorios privados dentro de la red de Walmart VPN."),
    ]

    for i, (icon, title, desc) in enumerate(benefits):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.7 + row * 2.7)
        add_rect(slide, x, y, Inches(4.0), Inches(2.4), WHITE)
        add_rect(slide, x, y, Inches(0.12), Inches(2.4), WALMART_YELLOW)
        add_text_box(slide, f"{icon}  {title}",
                     x + Inches(0.2), y + Inches(0.15),
                     Inches(3.7), Inches(0.55),
                     font_size=15, bold=True, color=WALMART_BLUE)
        add_text_box(slide, desc,
                     x + Inches(0.2), y + Inches(0.7),
                     Inches(3.7), Inches(1.5),
                     font_size=13, color=DARK_GRAY)


def slide_proximos_pasos(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_gradient(slide)
    slide_header(slide, "Próximos Pasos", "Plan de adopción del equipo")

    pasos = [
        ("Hoy",        WALMART_YELLOW,
         "Crear cuenta en github.com",
         "Solo toma 2 minutos. Usa tu correo Walmart o personal."),
        ("Esta semana", WHITE,
         "Conectar tu carpeta de trabajo",
         "Code Puppy te ayuda: 'Conecta mi carpeta a GitHub'."),
        ("Este mes",   WALMART_YELLOW,
         "Subir tus scripts actuales",
         "Primer commit con todos tus .py y .sql del proyecto."),
        ("Siempre",    WHITE,
         "Guardar cada cambio nuevo",
         "git add . && git commit && git push — o díselo al puppy 🐶"),
    ]

    for i, (tiempo, col, accion, detalle) in enumerate(pasos):
        y = Inches(1.6 + i * 1.4)
        # Línea de tiempo
        add_rect(slide, Inches(0.4), y + Inches(0.3),
                 Inches(12.5), Inches(0.05), WALMART_YELLOW)
        # Badge de tiempo
        add_rect(slide, Inches(0.4), y, Inches(1.6), Inches(0.55), WALMART_YELLOW)
        add_text_box(slide, tiempo,
                     Inches(0.4), y, Inches(1.6), Inches(0.55),
                     font_size=13, bold=True, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER)
        # Contenido
        add_rect(slide, Inches(2.2), y, Inches(10.6), Inches(1.1),
                 RGBColor(0x00, 0x3A, 0xA0))
        add_text_box(slide, accion,
                     Inches(2.4), y + Inches(0.05),
                     Inches(10.2), Inches(0.45),
                     font_size=15, bold=True, color=col)
        add_text_box(slide, detalle,
                     Inches(2.4), y + Inches(0.5),
                     Inches(10.2), Inches(0.5),
                     font_size=12, color=WHITE)

    # CTA final
    add_text_box(slide,
                 "¿Preguntas? Habla con Code Puppy 🐶  |  "
                 "puppy.walmart.com/doghouse",
                 Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
                 font_size=12, color=WALMART_YELLOW,
                 align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_portada(prs)
    slide_por_que(prs)
    slide_que_es(prs)
    slide_pasos(prs)
    slide_flujo_diario(prs)
    slide_seguridad(prs)
    slide_beneficios(prs)
    slide_proximos_pasos(prs)

    output = "Git_GitHub_Equipo_Walmart.pptx"
    prs.save(output)
    print(f"OK! Presentacion guardada: {output}")
    print(f"   Diapositivas: {len(prs.slides)}")


if __name__ == "__main__":
    main()
